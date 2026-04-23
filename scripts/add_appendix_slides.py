"""Add appendix/backup slides to the Digital Direction presentation.

These are reference slides for Q&A — not shown in the main flow, but available
if specific questions come up about methodology, golden data, merge logic, etc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Import helpers from main script
from build_presentation import (
    NAVY, BLUE, TEAL, WHITE, LIGHT_GRAY, DARK_GRAY, ORANGE, GREEN,
    set_slide_bg, add_text_box, add_bullet_list, add_box, add_card,
)

PURPLE = RGBColor(0xA8, 0x5C, 0xF8)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def add_appendix_header(slide, title, subtitle=""):
    """Standard appendix slide header with APPENDIX tag."""
    set_slide_bg(slide, NAVY)

    # APPENDIX tag
    tag = add_box(slide, 0.8, 0.3, 1.3, 0.35, DARK_GRAY)
    add_text_box(slide, 0.8, 0.3, 1.3, 0.35,
                 "APPENDIX", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, 2.3, 0.25, 10, 0.6,
                 title, font_size=32, color=WHITE, bold=True)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    if subtitle:
        add_text_box(slide, 0.8, 1.1, 11, 0.5, subtitle, font_size=16, color=LIGHT_GRAY)


def add_appendix_slides(prs):
    """Add all appendix slides to the presentation."""

    # ========================================
    # APPENDIX DIVIDER
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0, 2.5, 13.333, 1.0,
                 "APPENDIX", font_size=52, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0, 3.5, 13.333, 0.6,
                 "Reference slides for Q&A", font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # A: Golden Data — Clarity Needed from Client
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_appendix_header(slide, "Golden Data — Clarity Needed",
                        "Items where golden data expectations differ from what documents contain")

    gaps = [
        ("AT&T: Per-TN Address Mapping", RED, [
            "276 golden rows expect 16 different addresses",
            "CSR document contains only ONE address",
            "Analyst used an external master location DB",
            "Question: Can you provide the TN-to-address",
            "  mapping source as supplemental input?",
        ]),
        ("AT&T: Component Name Format", ORANGE, [
            "120 rows: accounts 5500/4339 use invoice format",
            '  e.g., "2.00 NETWORK ACCESS"',
            "Other accounts use CSR USOC descriptions",
            '  e.g., "Standard Centrex Feature"',
            "Question: Which format should be standard?",
        ]),
        ("Windstream: Circuit ID Format", ORANGE, [
            '801 rows use SDWAN IDs ("2389874-SDWAN-1")',
            "Our extraction uses CLLI format from invoices",
            '  e.g., "34/BBCW/201428/110/PUA /SDW"',
            "SDWAN IDs appear to be from a service portal",
            "Question: Source of SDWAN IDs?",
        ]),
    ]

    for i, (title, color, items) in enumerate(gaps):
        x = 0.5 + i * 4.2
        add_card(slide, x, 1.7, 3.8, 3.8, title, items, accent_color=color)

    # Bottom row: more gaps
    gaps2 = [
        ("AT&T: auto_renew Source", BLUE, [
            "457 rows missing (0% accuracy)",
            "Account 1586: golden says 'Yes' (251 rows)",
            "Accounts 5500/4339: golden says 'No'",
            "Not in CSR /TA or contract documents",
            "Question: Where does this come from?",
        ]),
        ("Golden 'Not Mentioned' Values", DARK_GRAY, [
            "contract_number: 205 rows = 'Not mentioned'",
            "contract_number: 44 rows = 'NA'",
            "billing_per_contract: 1,210 rows = 'N/A'",
            "These are analyst-entered placeholders",
            "Should we extract these or leave null?",
        ]),
        ("AT&T: Expiration Date Convention", BLUE, [
            "225 rows off by 1 day",
            "Our calc: begin + 12 months = Dec 2",
            "Golden expects: Dec 1 (day before)",
            "Convention: last day of term vs first after",
            "Question: Which convention to use?",
        ]),
    ]

    for i, (title, color, items) in enumerate(gaps2):
        x = 0.5 + i * 4.2
        add_card(slide, x, 5.6, 3.8, 1.8, title, items[:3], accent_color=color)

    # ========================================
    # B: Evaluation Methodology
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_appendix_header(slide, "Evaluation Methodology",
                        "How we measure and improve extraction accuracy")

    # Row matching
    add_card(slide, 0.5, 1.7, 3.8, 3.0, "Row Matching (3-Pass)", [
        "Pass 1: account + phone + charge_type + USOC",
        "  (most specific, lowest false match rate)",
        "Pass 2: account + phone (relaxed)",
        "  (handles missing charge_type/USOC)",
        "Pass 3: account + sub-account cross-match",
        "  (catches account numbering differences)",
        "Content pairing within groups by component name",
    ], accent_color=BLUE)

    # Field scoring
    add_card(slide, 4.7, 1.7, 3.8, 3.0, "Field Scoring", [
        "CORRECT: exact match (normalized)",
        "WRONG: different value extracted",
        "MISSING: golden has value, we don't",
        "EXTRA: we have value, golden doesn't",
        "Per-field: phone normalization (7/10 digit)",
        "Amount tolerance: +/- $0.01",
        "Date normalization across formats",
    ], accent_color=TEAL)

    # LLM Judge
    add_card(slide, 8.9, 1.7, 3.8, 3.0, "LLM Judge (Claude)", [
        "Evaluates fuzzy/semantic fields:",
        "  service_type, component_name, charge_type",
        "Telecom-aware scoring:",
        '  "BLC" = "Business Local Calling" \u2192 CORRECT',
        "Upgrades WRONG \u2192 CORRECT/PARTIAL only",
        "  (never downgrades deterministic scores)",
        "Batch evaluation: 20 fields per call",
    ], accent_color=PURPLE)

    # Categories
    add_text_box(slide, 0.5, 5.0, 12, 0.4,
                 "Accuracy Categories", font_size=18, color=WHITE, bold=True)

    cats = [
        ("Structured", BLUE, "Account #, phone #,\namounts, dates, zip", ">95%"),
        ("Semi-Structured", TEAL, "Address, billing name,\ncarrier name", ">85%"),
        ("Fuzzy", ORANGE, "Service type, component\nname, charge type", ">75%"),
        ("Contract", GREEN, "Term, dates, renewal,\nmonth-to-month", ">70%"),
    ]

    for i, (name, color, desc, target) in enumerate(cats):
        x = 0.5 + i * 3.15
        box = add_box(slide, x, 5.4, 2.9, 1.6, RGBColor(0x1A, 0x23, 0x3A), color)
        add_text_box(slide, x + 0.1, 5.5, 2.7, 0.35,
                     name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 5.85, 2.7, 0.7,
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 6.55, 2.7, 0.3,
                     f"Target: {target}", font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # ========================================
    # C: Smart Document Reading
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_appendix_header(slide, "Smart Document Reading",
                        "Three parsing paths, each optimized for different document types")

    paths = [
        ("Visual Path (Docling)", BLUE, [
            "For: invoices with complex tables,",
            "  scanned PDFs, two-column layouts",
            "Uses: Docling with TableFormer",
            "How: layout analysis \u2192 markdown export",
            "  preserving table structure",
            "Example: AT&T 14-page invoice with",
            "  nested sub-line charges",
        ]),
        ("Raw Text Path (pdfplumber)", TEAL, [
            "For: mainframe CSRs, simple PDFs",
            "  fixed-width column formats",
            "Uses: pdfplumber text extraction",
            "Chunking strategies:",
            "  \u2022 Page groups (3pp or 2pp for scanned)",
            "  \u2022 Section markers (---LISTINGS---)",
            "  \u2022 Boundary patterns (ACTIVITY FOR)",
            "Column-aware mode for multi-column",
        ]),
        ("Structured Path (pandas)", GREEN, [
            "For: CSV, XLSX, XLS spreadsheets",
            "Uses: pandas column mapping",
            "Direct schema mapping \u2014 no LLM needed",
            "Example: Peerless DID inventory (1500+",
            "  rows mapped instantly, $0 LLM cost)",
            "Windstream subscription exports",
        ]),
    ]

    for i, (title, color, items) in enumerate(paths):
        x = 0.5 + i * 4.2
        add_card(slide, x, 1.7, 3.8, 4.5, title, items, accent_color=color)

    # Global context injection
    add_text_box(slide, 0.5, 6.5, 12, 0.4,
                 "Global Context Injection", font_size=16, color=WHITE, bold=True)
    add_bullet_list(slide, 0.5, 6.9, 12, 0.6, [
        "SLA address lookup tables pre-extracted and injected into every chunk (AT&T CSR: per-TN addresses)",
        "Contract numbers (/CNUM) extracted once, available to all chunks for the same document",
    ], font_size=12, color=LIGHT_GRAY, bullet_color=TEAL, spacing=Pt(2))

    # ========================================
    # D: Cross-Document Merge
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_appendix_header(slide, "Cross-Document Merge",
                        "How we correlate data across invoices, CSRs, contracts, and reports")

    # Merge overview
    add_card(slide, 0.5, 1.7, 4.0, 2.5, "Priority Matrix", [
        "Each field has a priority per document type:",
        "  MRC/amounts: Invoice > Contract > CSR",
        "  Phone/BTN: CSR > Invoice > DID list",
        "  Address: CSR > Invoice (AT&T override)",
        "  Contract terms: Contract > CSR > Email",
        "Carrier-specific overrides in YAML config",
    ], accent_color=BLUE)

    add_card(slide, 4.9, 1.7, 4.0, 2.5, "Tiered Merge Keys", [
        "Tier 1: account + phone (strongest match)",
        "Tier 2: account + circuit (service-level)",
        "Tier 3: account only (enrichment only)",
        "Account equivalence: master \u2194 sub-account",
        "  linking via Union-Find algorithm",
        "Handles different numbering across docs",
    ], accent_color=TEAL)

    add_card(slide, 9.3, 1.7, 4.0, 2.5, "Conflict Resolution", [
        "Same-doc conflicts: fuller row wins",
        "Cross-doc conflicts: priority matrix decides",
        "True conflicts (equal priority, different values):",
        "  \u2192 Claude resolves with reasoning",
        "  \u2192 Picks based on data quality signals",
        "Source attribution tracked per field",
    ], accent_color=ORANGE)

    # Document roles
    add_text_box(slide, 0.5, 4.5, 12, 0.4,
                 "Document Roles in Merge", font_size=18, color=WHITE, bold=True)

    roles = [
        ("Primary", BLUE, [
            "Forms the merge base",
            "Invoices, CSRs",
            "Rows matched by Tier 1/2 keys",
            "Full field merge with priority",
        ]),
        ("Enrichment", TEAL, [
            "Fills gaps, never creates rows",
            "Contracts, service guides, emails",
            "Account-level fields propagated",
            "Contract terms enriched to services",
        ]),
        ("Supplemental", GREEN, [
            "Appended if no Tier 1/2 match",
            "DID lists, subscription exports",
            "Rows that don't match primary",
            "Preserves coverage",
        ]),
    ]

    for i, (title, color, items) in enumerate(roles):
        x = 0.8 + i * 4.1
        add_card(slide, x, 5.0, 3.7, 2.3, title, items, accent_color=color)

    # ========================================
    # E: LangFuse Observability
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_appendix_header(slide, "LangFuse: LLM Observability",
                        "Self-hosted tracing for every AI call in the pipeline")

    add_card(slide, 0.5, 1.7, 5.8, 3.0, "What We Track", [
        "Every Gemini extraction call: prompt, response, tokens, latency",
        "Every Claude merge/eval call: conflict resolution reasoning",
        "Per-document cost breakdown (input + output tokens \u00d7 model rate)",
        "Extraction quality signals: confidence scores, validation flags",
        "Correction-to-diagnosis traces: which prompt produced which error",
        "Token usage trends across carriers and document types",
    ], accent_color=BLUE)

    add_card(slide, 6.7, 1.7, 5.8, 3.0, "Why It Matters", [
        "Debug extraction issues: see exactly what the AI read and produced",
        "Optimize prompts: identify which instructions cause errors",
        "Track cost: per-carrier, per-doc-type cost visibility",
        "Audit trail: every extraction decision is traceable",
        "Self-healing visibility: see how corrections change future outputs",
        "No external SaaS: self-hosted alongside PostgreSQL in Docker",
    ], accent_color=TEAL)

    # Architecture diagram (text-based)
    add_text_box(slide, 0.5, 5.0, 12, 0.4,
                 "Integration Architecture", font_size=18, color=WHITE, bold=True)

    arch_box = add_box(slide, 0.5, 5.4, 12.3, 1.8, RGBColor(0x12, 0x1A, 0x2E), DARK_GRAY)
    add_text_box(slide, 0.8, 5.5, 11.7, 1.6,
                 "Upload \u2192 Classify \u2192 Parse \u2192 Extract (Gemini) \u2192 Merge (Claude) \u2192 Validate \u2192 Review UI\n"
                 "                                        \u2502                       \u2502\n"
                 "                                   LangFuse Trace          LangFuse Trace\n"
                 "                                        \u2502                       \u2502\n"
                 "                                  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500 LangFuse Dashboard (localhost:3100) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
                 "                                        Prompt \u2022 Response \u2022 Tokens \u2022 Cost \u2022 Latency",
                 font_size=11, color=LIGHT_GRAY, font_name="Courier New")

    return prs


def main():
    # Load existing presentation
    prs = Presentation("/Users/techjays/Desktop/Digital_Direction_Platform_Overview.pptx")

    # Add appendix slides
    add_appendix_slides(prs)

    # Save
    output_path = "/Users/techjays/Desktop/Digital_Direction_Platform_Overview.pptx"
    prs.save(output_path)
    print(f"Updated presentation with appendix: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
