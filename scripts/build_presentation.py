"""Build the Digital Direction client presentation (PPTX).

Addresses Matt Hammer's email concerns:
1. Does the platform need exact template matches? (NO)
2. Can we provide inputs only for a third customer? (YES)
3. Contract intelligence — month-to-month, contract correlation, billing compliance
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy background
BLUE = RGBColor(0x1E, 0x90, 0xFF)       # Primary blue accent
TEAL = RGBColor(0x00, 0xC9, 0xA7)       # Teal accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)     # Warning/highlight
GREEN = RGBColor(0x00, 0xC8, 0x53)      # Success
SOFT_BG = RGBColor(0xF5, 0xF7, 0xFA)    # Light background


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=TEAL, font_name="Calibri", spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "  # Small right triangle
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name

        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name

    return txBox


def add_box(slide, left, top, width, height, fill_color, border_color=None, radius=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items,
             accent_color=BLUE, bg_color=RGBColor(0x1A, 0x23, 0x3A)):
    """Add a styled card with title and bullet items."""
    # Card background
    add_box(slide, left, top, width, height, bg_color, accent_color)

    # Accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    add_text_box(slide, left + 0.2, top + 0.15, width - 0.4, 0.4,
                 title, font_size=14, color=accent_color, bold=True)

    # Body items
    if body_items:
        add_bullet_list(slide, left + 0.15, top + 0.55, width - 0.3, height - 0.7,
                        body_items, font_size=11, color=LIGHT_GRAY, bullet_color=accent_color,
                        spacing=Pt(4))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========================================
    # SLIDE 1: Title
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.0), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 1, 2.2, 11, 1.2,
                 "Digital Direction", font_size=52, color=WHITE, bold=True)
    add_text_box(slide, 1, 3.4, 11, 0.8,
                 "AI-Powered Telecom Document Intelligence Platform",
                 font_size=24, color=TEAL)
    add_text_box(slide, 1, 4.5, 11, 0.5,
                 "Platform Overview & Methodology",
                 font_size=18, color=DARK_GRAY)
    add_text_box(slide, 1, 6.5, 5, 0.4,
                 "Presented by TechJays", font_size=14, color=DARK_GRAY)
    add_text_box(slide, 7, 6.5, 5, 0.4,
                 "Confidential", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)

    # ========================================
    # SLIDE 2: The Challenge
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "The Challenge", font_size=36, color=WHITE, bold=True)

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.6,
                 "Telecom inventory management is manual, error-prone, and doesn't scale.",
                 font_size=18, color=LIGHT_GRAY)

    # Three problem cards
    problems = [
        ("Document Overload", [
            "Invoices, CSRs, contracts, service orders",
            "Each carrier uses different formats",
            "100s of pages per enterprise account",
            "New formats appear without warning",
        ]),
        ("Manual Data Entry", [
            "Analysts manually copy 60+ fields per row",
            "Cross-referencing across 3-5 documents",
            "Error-prone with high-density data",
            "Days of work per client inventory",
        ]),
        ("Contract Complexity", [
            "Matching contracts to active services",
            "Detecting month-to-month status",
            "Verifying billing matches contract terms",
            "Tracking renewals and expirations",
        ]),
    ]

    for i, (title, items) in enumerate(problems):
        x = 0.8 + i * 4.1
        add_card(slide, x, 2.2, 3.7, 4.5, title, items, accent_color=ORANGE)

    add_text_box(slide, 0.8, 7.0, 11, 0.4,
                 "Result: weeks of analyst time per client, inconsistent outputs, no scalability path.",
                 font_size=14, color=ORANGE, bold=True)

    # ========================================
    # SLIDE 3: Our Approach (KEY — addresses "template matching" question)
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Our Approach: Zero-Template Intelligence", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Key differentiator callout
    callout_bg = add_box(slide, 0.8, 1.4, 11.7, 1.2, RGBColor(0x0A, 0x2A, 0x1A), TEAL)
    add_text_box(slide, 1.1, 1.5, 11, 0.4,
                 "KEY DIFFERENTIATOR", font_size=12, color=TEAL, bold=True)
    add_text_box(slide, 1.1, 1.85, 11, 0.7,
                 "No template matching required. The platform understands telecom documents through "
                 "carrier domain knowledge and AI reasoning, not rigid templates. "
                 "Give it a new carrier's invoice, and it extracts structured data on the first attempt.",
                 font_size=16, color=WHITE)

    # How it works - 3 pillars
    pillars = [
        ("Carrier Knowledge", BLUE, [
            "Understands AT&T, Windstream, Spectrum,",
            "Peerless document structures",
            "USOC codes, field codes, service types",
            "Account number formats per carrier",
            "New carriers: config-driven onboarding",
        ]),
        ("AI Reasoning", TEAL, [
            "Gemini extracts meaning, not patterns",
            "Handles format variations naturally",
            "OCR + multimodal for scanned docs",
            "Cross-document correlation",
            "No training data needed per template",
        ]),
        ("Self-Healing", GREEN, [
            "Human corrections feed back into system",
            "Root-cause diagnosis per correction",
            "Accuracy improves with every review",
            "Domain knowledge grows automatically",
            "Reduces manual review over time",
        ]),
    ]

    for i, (title, color, items) in enumerate(pillars):
        x = 0.8 + i * 4.1
        add_card(slide, x, 2.9, 3.7, 4.0, title, items, accent_color=color)

    # ========================================
    # SLIDE 4: Architecture Overview
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Platform Architecture", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.5,
                 "5-Stage Pipeline: Documents In, Structured Inventory Out",
                 font_size=18, color=LIGHT_GRAY)

    # Pipeline stages as connected boxes
    stages = [
        ("1. CLASSIFY", "Identify carrier,\ndoc type, format", BLUE),
        ("2. PARSE", "OCR, text extraction,\nchunking", RGBColor(0x6C, 0x5C, 0xE7)),
        ("3. EXTRACT", "AI extracts 60 fields\nper row", TEAL),
        ("4. MERGE", "Cross-doc correlation,\npriority rules", ORANGE),
        ("5. VALIDATE", "Confidence scoring,\nquality checks", GREEN),
    ]

    for i, (title, desc, color) in enumerate(stages):
        x = 0.5 + i * 2.5
        y = 2.1

        # Stage box
        box = add_box(slide, x, y, 2.2, 1.6, RGBColor(0x1A, 0x23, 0x3A), color)

        # Stage title
        add_text_box(slide, x + 0.1, y + 0.15, 2.0, 0.35,
                     title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Stage description
        add_text_box(slide, x + 0.1, y + 0.55, 2.0, 0.9,
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between stages (except last)
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.25), Inches(y + 0.6),
                Inches(0.25), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # Bottom section: what feeds into the pipeline
    add_text_box(slide, 0.8, 4.2, 11, 0.5,
                 "What Goes In", font_size=20, color=WHITE, bold=True)

    inputs = [
        ("Documents", BLUE, ["Invoices (PDF, scanned)", "CSRs (mainframe text)", "Contracts (PDF)",
                              "Service orders, emails"]),
        ("Carrier Config", TEAL, ["YAML-driven, not code", "Format signatures", "USOC code mappings",
                                   "Merge priority rules"]),
        ("Corrections", GREEN, ["Human review feedback", "Root-cause diagnosed", "Routes to correct layer",
                                 "Builds domain knowledge"]),
    ]

    for i, (title, color, items) in enumerate(inputs):
        x = 0.8 + i * 4.1
        add_card(slide, x, 4.8, 3.7, 2.4, title, items, accent_color=color)

    # ========================================
    # SLIDE 5: Document Intelligence
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Any Document, Any Carrier", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.6,
                 "The platform handles documents through understanding, not memorization.",
                 font_size=18, color=LIGHT_GRAY)

    # Two columns: How it works vs Traditional
    # Left: Our approach
    add_card(slide, 0.8, 2.2, 5.6, 4.8,
             "Digital Direction Approach", [
                 "AI reads the document like an analyst would",
                 "Carrier domain knowledge provides context (not templates)",
                 "Format detection identifies document structure automatically",
                 "Same carrier, different format? Handled seamlessly",
                 "New carrier? Add a YAML config file, extract immediately",
                 "Scanned PDFs: multimodal AI reads the image directly",
                 "Cross-document merge fills gaps from multiple sources",
                 "No training required on specific invoice layouts",
             ], accent_color=TEAL)

    # Right: Traditional approach
    add_card(slide, 6.8, 2.2, 5.6, 4.8,
             "Traditional Template Approach", [
                 "Requires exact template match per invoice format",
                 "New format = new template build (days/weeks)",
                 "Template breaks when carrier changes layout",
                 "Can't handle format variations within same carrier",
                 "Scanned PDFs often unsupported",
                 "No cross-document correlation",
                 "Manual maintenance as formats evolve",
                 "Scale = more templates = more maintenance",
             ], accent_color=ORANGE)

    # ========================================
    # SLIDE 6: Self-Healing System
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Self-Healing Intelligence", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.6,
                 "Every correction makes the system smarter. The platform diagnoses WHY it was wrong, "
                 "not just WHAT was wrong.",
                 font_size=18, color=LIGHT_GRAY)

    # Self-healing flow as horizontal steps
    steps = [
        ("Reviewer Corrects", "Analyst fixes a field\nin the review UI", BLUE),
        ("System Diagnoses", "Traces error through\npipeline stages", TEAL),
        ("Root Cause Found", "Extraction? Merge?\nData gap? Enrichment?", ORANGE),
        ("Fix Applied", "Routes fix to the\ncorrect layer", GREEN),
        ("System Learns", "Next extraction avoids\nthe same mistake", RGBColor(0xA8, 0x5C, 0xF8)),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = 0.4 + i * 2.5
        y = 2.2

        box = add_box(slide, x, y, 2.2, 1.8, RGBColor(0x1A, 0x23, 0x3A), color)
        add_text_box(slide, x + 0.1, y + 0.15, 2.0, 0.4,
                     title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, y + 0.6, 2.0, 1.0,
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.25), Inches(y + 0.65),
                Inches(0.25), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # Root cause breakdown
    add_text_box(slide, 0.8, 4.4, 11, 0.5,
                 "Root-Cause Diagnosis Categories", font_size=20, color=WHITE, bold=True)

    causes = [
        ("EXTRACTION", "AI misread the\nsource document", "Fix: Prompt improvement\nor domain knowledge", BLUE),
        ("MERGE", "Wrong priority in\ncross-doc merge", "Fix: Adjust merge rules\nfor this carrier", TEAL),
        ("ENRICHMENT", "Cross-doc data\nnot propagated", "Fix: Expand enrichment\nlookup scope", ORANGE),
        ("DATA GAP", "Value not in any\nsource document", "Action: Request\nsupplemental data", RGBColor(0xE7, 0x4C, 0x3C)),
    ]

    for i, (title, desc, fix, color) in enumerate(causes):
        x = 0.8 + i * 3.1
        add_card(slide, x, 5.0, 2.8, 2.2, title, [desc, fix], accent_color=color)

    # ========================================
    # SLIDE 7: Contract Intelligence
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Contract Intelligence", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.6,
                 "Cross-referencing invoices, CSRs, and contracts to build a complete inventory picture.",
                 font_size=18, color=LIGHT_GRAY)

    # Three capability cards addressing Matt's specific asks
    capabilities = [
        ("Month-to-Month Detection", BLUE, [
            "Extracts contract dates from CSR /TA fields",
            "Extracts terms from contract documents",
            "Derives M2M status: expired contract = M2M",
            "Propagates contract status across all services",
            "Automatically flags expiring contracts",
        ]),
        ("Contract-Service Correlation", TEAL, [
            "Matches contracts to services via account/TN",
            "Cross-doc merge links invoice + CSR + contract",
            "Contract number extracted from /CNUM fields",
            "Term months, begin/end dates correlated",
            "Handles multi-service contracts",
        ]),
        ("Billing Compliance", GREEN, [
            "Compares contracted rates vs billed amounts",
            "Detects overcharges and billing errors",
            "Flags services billed without contracts",
            "Identifies discount expiration impacts",
            "Generates compliance audit trail",
        ]),
    ]

    for i, (title, color, items) in enumerate(capabilities):
        x = 0.8 + i * 4.1
        add_card(slide, x, 2.2, 3.7, 4.5, title, items, accent_color=color)

    # Bottom note
    add_text_box(slide, 0.8, 7.0, 11, 0.4,
                 "Contract intelligence is built into the extraction pipeline, not a separate system.",
                 font_size=14, color=TEAL, bold=True)

    # ========================================
    # SLIDE 8: New Customer Onboarding
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "New Customer Onboarding", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Directly addressing: "Can we provide inputs only for a third customer?"
    callout_bg = add_box(slide, 0.8, 1.4, 11.7, 1.0, RGBColor(0x0A, 0x2A, 0x1A), TEAL)
    add_text_box(slide, 1.1, 1.5, 11, 0.8,
                 'Yes. Provide documents only. No golden data, no template matching, no prior examples needed. '
                 'If the carrier is already supported, extraction works immediately. '
                 'If it\'s a new carrier, a carrier config is added (typically < 1 day).',
                 font_size=16, color=WHITE)

    # Timeline
    add_text_box(slide, 0.8, 2.8, 11, 0.5,
                 "What Onboarding Looks Like", font_size=20, color=WHITE, bold=True)

    timeline = [
        ("Day 1", "Upload Documents", "Upload invoices, CSRs, contracts.\nPlatform classifies automatically.", BLUE),
        ("Day 1", "First Extraction", "AI extracts 60 fields per row.\nReview results in the web UI.", TEAL),
        ("Day 2-3", "Human Review", "Analyst reviews, corrects errors.\nSystem diagnoses root causes.", ORANGE),
        ("Day 3-5", "Accuracy Climbs", "Self-healing improves extraction.\nDomain knowledge grows.", GREEN),
        ("Ongoing", "Production", "Corrections decrease over time.\nNew docs processed automatically.", RGBColor(0xA8, 0x5C, 0xF8)),
    ]

    for i, (time, title, desc, color) in enumerate(timeline):
        x = 0.4 + i * 2.5
        y = 3.4

        box = add_box(slide, x, y, 2.2, 2.8, RGBColor(0x1A, 0x23, 0x3A), color)

        # Time label
        add_text_box(slide, x + 0.1, y + 0.1, 2.0, 0.3,
                     time, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + 0.1, y + 0.4, 2.0, 0.4,
                     title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + 0.1, y + 0.9, 2.0, 1.6,
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # What we need from client
    add_text_box(slide, 0.8, 6.5, 11, 0.4,
                 "What We Need From You", font_size=18, color=WHITE, bold=True)

    add_bullet_list(slide, 0.8, 6.9, 11, 0.6, [
        "Source documents (invoices, CSRs, contracts) for the client's carriers",
        "Carrier list (so we can verify config coverage or add new carriers)",
    ], font_size=14, color=LIGHT_GRAY, bullet_color=TEAL, spacing=Pt(4))

    # ========================================
    # SLIDE 9: Supported Carriers & Documents
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Carrier & Document Coverage", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Carrier cards
    carriers = [
        ("AT&T", BLUE, [
            "Invoices (multi-page, two-column)",
            "CSRs (box format + section-marker)",
            "Contracts (service agreements)",
            "Service guides",
            "USOC code resolution",
        ]),
        ("Windstream", TEAL, [
            "Enterprise invoices (470+ pages)",
            "Kinetic invoices (small business)",
            "CSRs and subscription exports",
            "Contracts (SDWAN, fiber)",
            "Sub-account level extraction",
        ]),
        ("Spectrum", ORANGE, [
            "Consolidated multi-location bills",
            "Enterprise fiber invoices",
            "Scanned PDFs (multimodal OCR)",
            "Contracts and service orders",
            "Account hierarchy detection",
        ]),
        ("Peerless Network", GREEN, [
            "SIP trunk invoices",
            "DID inventory (CSV/XLSX)",
            "Subscription exports",
            "Quotes with contract terms",
            "Channel fee + per-DID pricing",
        ]),
    ]

    for i, (name, color, items) in enumerate(carriers):
        x = 0.5 + i * 3.15
        add_card(slide, x, 1.5, 2.9, 4.0, name, items, accent_color=color)

    # Document types
    add_text_box(slide, 0.8, 5.8, 11, 0.5,
                 "Document Types Supported", font_size=18, color=WHITE, bold=True)

    doc_types = [
        "PDF Invoices (text & scanned)", "CSRs (mainframe text)", "Contracts (PDF)",
        "CSV/XLSX Reports", "Email (.msg/.eml)", "Service Orders", "DOCX Agreements",
    ]

    add_bullet_list(slide, 0.8, 6.3, 11, 1.0, doc_types,
                    font_size=13, color=LIGHT_GRAY, bullet_color=TEAL, spacing=Pt(2))

    # ========================================
    # SLIDE 10: The 60-Field Output
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Standardized 60-Field Output", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_text_box(slide, 0.8, 1.3, 11, 0.5,
                 "Every carrier, every document type, normalized into one consistent schema.",
                 font_size=18, color=LIGHT_GRAY)

    # Field categories
    categories = [
        ("Identity", BLUE, [
            "Carrier name & account numbers",
            "Master/sub-account hierarchy",
            "BTN & phone numbers",
            "Circuit IDs",
        ]),
        ("Location", TEAL, [
            "Service address (street, city, state, zip)",
            "Z-location for circuits",
            "Billing name",
            "Country & currency",
        ]),
        ("Service", ORANGE, [
            "Service type (POTS, Centrex, SIP, etc.)",
            "USOC codes & component names",
            "MRC, quantity, cost per unit",
            "Charge type (MRC, usage, tax)",
        ]),
        ("Contract", GREEN, [
            "Contract number & term months",
            "Begin & expiration dates",
            "Auto-renew status",
            "Month-to-month detection",
        ]),
    ]

    for i, (title, color, items) in enumerate(categories):
        x = 0.8 + i * 3.1
        add_card(slide, x, 2.0, 2.8, 3.2, title, items, accent_color=color)

    # Output format
    add_text_box(slide, 0.8, 5.5, 11, 0.5,
                 "Output Formats", font_size=18, color=WHITE, bold=True)

    add_bullet_list(slide, 0.8, 6.0, 11, 1.2, [
        "Excel export with confidence color-coding per cell (green/yellow/red)",
        "JSON API for integration with inventory management systems",
        "Web UI for human review with inline editing and bulk approve",
        "Corrections import: upload corrected Excel, system learns from differences",
    ], font_size=14, color=LIGHT_GRAY, bullet_color=TEAL, spacing=Pt(6))

    # ========================================
    # SLIDE 11: Addressing Your Questions
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Addressing Your Questions", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Q&A format
    qa_pairs = [
        ('"Is a previous exact template match required for every possible invoice?"',
         "No. The platform uses carrier domain knowledge + AI reasoning, not template matching. "
         "A new invoice format from a known carrier (AT&T, Windstream, etc.) is processed "
         "immediately. Format variations are handled by the AI, not by rigid templates."),
        ('"Can we provide inputs only for a third customer?"',
         "Yes. Upload the customer's documents. If the carriers are already configured, "
         "extraction runs immediately. If new carriers are involved, we add a carrier config "
         "(typically under 1 day). No golden data or prior outputs needed."),
        ('"Can the platform determine month-to-month status, correlate contracts, and verify billing?"',
         "Yes. The platform extracts contract terms from CSRs and contract documents, "
         "derives month-to-month status from expiration dates, and correlates contract "
         "numbers to specific services. Billing compliance (contracted vs actual rates) "
         "is built into the cross-document merge layer."),
    ]

    y = 1.5
    for q, a in qa_pairs:
        # Question
        add_text_box(slide, 1.0, y, 11.3, 0.5,
                     q, font_size=14, color=TEAL, bold=True)
        # Answer
        add_text_box(slide, 1.0, y + 0.45, 11.3, 1.0,
                     a, font_size=14, color=LIGHT_GRAY)
        y += 1.7

    # ========================================
    # SLIDE 12: Next Steps
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, 0.8, 0.4, 11, 0.6,
                 "Next Steps", font_size=36, color=WHITE, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    steps = [
        ("1", "Third Customer Documents", "Provide invoices, CSRs, and contracts for the third customer. "
         "We'll process them through the platform and deliver initial extraction results.",
         BLUE),
        ("2", "Contract Intelligence Demo", "Live demonstration of contract-service correlation, "
         "month-to-month detection, and billing compliance checking on your data.",
         TEAL),
        ("3", "Review & Feedback Cycle", "Your analysts review the extracted data in our web UI. "
         "Every correction improves the system through self-healing feedback.",
         ORANGE),
        ("4", "Production Readiness", "Once accuracy meets targets, move to production deployment "
         "with ongoing self-improvement and new carrier onboarding.",
         GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        y = 1.5 + i * 1.4

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.1), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text_box(slide, 1.0, y + 0.12, 0.5, 0.5,
                     num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Title + description
        add_text_box(slide, 1.8, y, 10, 0.4,
                     title, font_size=18, color=WHITE, bold=True)
        add_text_box(slide, 1.8, y + 0.4, 10, 0.8,
                     desc, font_size=14, color=LIGHT_GRAY)

    # Contact
    add_text_box(slide, 0.8, 7.0, 11, 0.4,
                 "Contact: rajat@techjays.com  |  TechJays",
                 font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # Save
    # ========================================
    output_path = "Digital_Direction_Platform_Overview.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
